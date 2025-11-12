import os
import io
import csv
import tempfile
from pathlib import Path
import unicodedata

import streamlit as st
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import requests

# --------------------------------------------------------------------------------------
# NORMALIZATION + PATH HELPERS (define FIRST so other functions can use them safely)
# --------------------------------------------------------------------------------------
def _norm(s: str) -> str:
    """Simple normalization: lowercase, strip, collapse spaces, NFKC unicode."""
    s = str(s or "")
    s = unicodedata.normalize("NFKC", s)
    return " ".join(s.strip().split()).lower()

def _tokens(s: str) -> set:
    """Tokenize a type string into words for tolerant matching ('coffee table' -> {'coffee','table'})."""
    s = _norm(s).replace("-", " ").replace("_", " ")
    return set([t for t in s.split() if t])

def _child_caseless(parent: Path, wanted: str) -> Path | None:
    wanted_n = _norm(wanted)
    if not parent.exists() or not parent.is_dir():
        return None
    for p in parent.iterdir():
        try:
            if p.is_dir() and _norm(p.name) == wanted_n:
                return p
        except Exception:
            continue
    return None

def resolve_caseless_path(base_dir: str | Path, *segments: str) -> Path | None:
    cur = Path(base_dir)
    for seg in segments:
        nxt = _child_caseless(cur, seg)
        if nxt is None:
            return None
        cur = nxt
    return cur

def show_image_safe(src):
    try:
        s = str(src).strip()
        if not s:
            st.caption("⚠️ Empty image reference")
            return
        if "://" in s:
            r = requests.get(s, timeout=30)
            r.raise_for_status()
            img = Image.open(io.BytesIO(r.content)).convert("RGB")
        else:
            img = Image.open(s).convert("RGB")
        st.image(img, use_column_width=True)
    except Exception as e:
        st.caption(f"⚠️ Failed to preview image ({src}): {e}")

# --------------------------------------------------------------------------------------
# FILE/ASSET CONSTANTS
# --------------------------------------------------------------------------------------
BASE_DIR = Path(__file__).parent
EXCEL_PATH = BASE_DIR / "all companys database.xlsx"
IMAGE_BASE = BASE_DIR / "images"
LOGO_BASE = BASE_DIR / "static" / "logo"
FIRST_PATH = BASE_DIR / "static" / "img" / "first.png"
LAST_PATH  = BASE_DIR / "static" / "img" / "last.png"
LOCAL_MANIFEST = BASE_DIR / "image_manifest.csv"

# --------------------------------------------------------------------------------------
# DATA LOADERS
# --------------------------------------------------------------------------------------
@st.cache_data(show_spinner=False)
def load_excel(path: Path) -> pd.DataFrame:
    return pd.read_excel(path)

@st.cache_data(show_spinner=False)
def load_manifest():
    """
    Manifest CSV columns:
      Company, Product, Type, ImageURLs
    ImageURLs is '|' separated list of absolute URLs.
    Also registers a swapped (Product<->Type) key to tolerate mismatches
    between your Excel and CSV layouts.
    """
    csv_text = st.secrets.get("IMAGE_MANIFEST_CSV", "")
    if csv_text:
        rows = list(csv.DictReader(csv_text.splitlines()))
    rows = []
    try:
        if url:
            resp = requests.get(url, timeout=30)
            resp.raise_for_status()
            rows = list(csv.DictReader(resp.text.splitlines()))
        elif LOCAL_MANIFEST.exists():
            with open(LOCAL_MANIFEST, newline="", encoding="utf-8") as f:
                rows = list(csv.DictReader(f))
    except Exception as e:
        st.warning(f"Could not load image manifest: {e}")

    def _clean_url(u: str) -> str:
        u = (u or "").strip()
        # strip accidental trailing slash after extension
        if u.lower().endswith((".jpg/", ".jpeg/", ".png/", ".webp/")):
            u = u[:-1]
        return u

    manifest = {}
    for r in rows:
        c_raw, p_raw, t_raw = r.get("Company", ""), r.get("Product", ""), r.get("Type", "")
        c, p, t = _norm(c_raw), _norm(p_raw), _norm(t_raw)
        urls = [ _clean_url(u) for u in (r.get("ImageURLs") or "").split("|") if _clean_url(u) ]
        if c and p and t and urls:
            # main key
            manifest[(c, p, t)] = urls
            # tolerant swapped key (handles CSV vs Excel column mismatch)
            manifest[(c, t, p)] = urls
    return manifest


DATA = load_excel(EXCEL_PATH)
MANIFEST = load_manifest()

# --------------------------------------------------------------------------------------
# IMAGE RESOLUTION (CSV/URL first, local images/ fallback)
# --------------------------------------------------------------------------------------
def get_image_list(company: str, product: str, ptype: str):
    """Return image list for (Company, Product, Type), tolerant to Product/Type swap."""
    c, p, t = _norm(company), _norm(product), _norm(ptype)

    def _match_from_manifest(c, p, t):
        # 1) Exact (includes swapped key thanks to load_manifest)
        if (c, p, t) in MANIFEST:
            return MANIFEST[(c, p, t)]

        # 2) soft: startswith / contains (same c,p)
        for (mc, mp, mt), urls in MANIFEST.items():
            if mc == c and mp == p and (mt == t or mt.startswith(t) or t in mt):
                return urls

        # 3) token-based partial (same c,p)
        t_tokens = _tokens(t)
        best = None
        best_overlap = 0
        for (mc, mp, mt), urls in MANIFEST.items():
            if mc == c and mp == p:
                mt_tokens = _tokens(mt)
                overlap = len(t_tokens & mt_tokens)
                if overlap > best_overlap:
                    best_overlap = overlap
                    best = urls
        if best and best_overlap > 0:
            return best

        # 4) fallback: any (c,p)
        for (mc, mp, mt), urls in MANIFEST.items():
            if mc == c and mp == p:
                return urls
        return []

    # Try normal ordering
    urls = _match_from_manifest(c, p, t)
    if urls:
        return urls

    # As an extra safety net, try reading Type/Product swapped (helps when Excel is flipped)
    urls = _match_from_manifest(c, t, p)
    if urls:
        return urls

    # Local filesystem fallback (case-insensitive): images/Company/Product/Type/*
    folder = resolve_caseless_path(IMAGE_BASE, company, product, ptype)
    images = []
    if folder and folder.exists():
        for file in sorted(folder.iterdir()):
            if file.suffix.lower() in (".jpg", ".jpeg", ".png", ".webp"):
                images.append(str(file))
    return images

# --------------------------------------------------------------------------------------
# PPT CREATION (unchanged styling)
# --------------------------------------------------------------------------------------
def get_scaled_dimensions(img, max_width, max_height):
    img_width_px, img_height_px = img.size
    aspect_ratio = img_width_px / img_height_px
    box_aspect = max_width / max_height
    if aspect_ratio > box_aspect:
        width = max_width
        height = width / aspect_ratio
    else:
        height = max_height
        width = height * aspect_ratio
    return width, height

def open_pil_image(path_or_url):
    if "://" in str(path_or_url):
        resp = requests.get(path_or_url, timeout=60)
        resp.raise_for_status()
        return Image.open(io.BytesIO(resp.content))
    return Image.open(path_or_url)

def fetch_to_tempfile(path_or_url):
    if "://" not in str(path_or_url):
        return path_or_url
    resp = requests.get(path_or_url, timeout=60)
    resp.raise_for_status()
    ext = ".png" if str(path_or_url).lower().endswith(".png") else ".jpg"
    fd, tpath = tempfile.mkstemp(suffix=ext)
    with os.fdopen(fd, "wb") as f:
        f.write(resp.content)
    return tpath

def create_beautiful_ppt(slide_data_list, include_intro_outro=True):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    if include_intro_outro and FIRST_PATH.exists():
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(FIRST_PATH), Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)

    for slide_data in slide_data_list:
        slide = prs.slides.add_slide(blank)
        company = slide_data.get('company', '')
        product = slide_data.get('product', '')
        link = slide_data.get('link', '')
        imgs = slide_data.get('images', [])

        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12), Inches(0.7))
        frame = title_shape.text_frame
        frame.text = product
        p = frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        # Images grid
        y_img_top = 1.2
        y_img_bottom = 6.9
        max_img_height = y_img_bottom - y_img_top
        slide_width_in = prs.slide_width.inches

        img_count = len(imgs)
        if img_count > 0:
            padding = 0.2
            columns = min(img_count, 3)
            rows = (img_count + columns - 1) // columns
            available_width = slide_width_in - (padding * (columns + 1))
            available_height = max_img_height - ((rows - 1) * padding)
            cell_width = available_width / columns
            cell_height = available_height / rows

            for i, img_src in enumerate(imgs):
                row = i // columns
                col = i % columns
                try:
                    with open_pil_image(img_src) as img:
                        img_width, img_height = get_scaled_dimensions(img, max_width=cell_width, max_height=cell_height)
                except Exception:
                    img_width, img_height = cell_width, cell_height

                x = padding + col * (cell_width + padding) + (cell_width - img_width) / 2
                y = y_img_top + row * (cell_height + padding) + (cell_height - img_height) / 2
                add_path = fetch_to_tempfile(img_src)
                slide.shapes.add_picture(add_path, Inches(x), Inches(y),
                                         width=Inches(img_width), height=Inches(img_height))

        # Logo
        logo_dir = resolve_caseless_path(LOGO_BASE, company)
        logo_path = None
        if logo_dir and logo_dir.exists():
            for ext in (".png", ".jpg", ".jpeg"):
                cand = logo_dir / f"logo{ext}"
                if cand.exists():
                    logo_path = str(cand)
                    break
        if logo_path:
            slide.shapes.add_picture(logo_path, prs.slide_width - Inches(1.2), Inches(0.1), width=Inches(1.1))

        # Footer
        cp_box = slide.shapes.add_textbox(prs.slide_width - Inches(3.6), prs.slide_height - Inches(0.3), Inches(3.6), Inches(0.4))
        cp_frame = cp_box.text_frame
        cp_frame.text = "Copyright © 2025 Altossa Projects LLp. All Rights Reserved."
        cp_para = cp_frame.paragraphs[0]
        cp_para.font.size = Pt(10)
        cp_para.font.color.rgb = RGBColor(128, 128, 128)

        if link:
            link_box = slide.shapes.add_textbox(Inches(0.1), prs.slide_height - Inches(0.3), Inches(7), Inches(0.4))
            link_frame = link_box.text_frame
            link_frame.text = str(link)
            p = link_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0, 102, 204)

    if include_intro_outro and LAST_PATH.exists():
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(LAST_PATH), Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)

    return prs

# --------------------------------------------------------------------------------------
# SESSION STATE
# --------------------------------------------------------------------------------------
if 'ppt_items' not in st.session_state:
    st.session_state.ppt_items = {}
if 'temp_selection' not in st.session_state:
    st.session_state.temp_selection = {}
if 'last_temp_key' not in st.session_state:
    st.session_state.last_temp_key = None
if 'ppt_ready' not in st.session_state:
    st.session_state.ppt_ready = False
if 'ppt_path' not in st.session_state:
    st.session_state.ppt_path = None

# --------------------------------------------------------------------------------------
# UI (kept same behavior)
# --------------------------------------------------------------------------------------
st.title("Product Selector with Search")

# Dev/status hints (won't affect UX)
with st.expander("Status (debug tips)", expanded=False):
    st.write(f"Excel loaded: `{EXCEL_PATH.name}` → {len(DATA)} rows")
    st.write(f"Manifest source: {'st.secrets URL' if st.secrets.get('IMAGE_MANIFEST_URL') else ('image_manifest.csv' if LOCAL_MANIFEST.exists() else '— none —')}")
    if MANIFEST:
        st.write(f"Manifest keys: {len(MANIFEST)}")

search_query = st.text_input("Search by Type", "")

if 'search_selection_keys' not in st.session_state:
    st.session_state.search_selection_keys = set()

if search_query:
    filtered_data = DATA[DATA['Type'].str.contains(search_query, case=False, na=False)]
    for idx, row in filtered_data.iterrows():
        company, product, ptype, link = row['Company'], row['Product'], row['Type'], row.get('Link', '')
        img_paths = [p for p in get_image_list(company, product, ptype) if p and str(p).strip()]

        st.markdown(f"### {product} - {ptype}")
        if img_paths:
            cols = st.columns(min(4, len(img_paths)))
            selected_imgs = []
            for i, path in enumerate(img_paths):
                with cols[i % len(cols)]:
                    show_image_safe(path)
                    key = f"search_{company}_{product}_{ptype}_{i}".replace(" ", "_")
                    if st.checkbox("Include", key=key):
                        selected_imgs.append(path)
            if selected_imgs:
                st.session_state.ppt_items[f"{company}_{product}_{ptype}".replace(" ", "_")] = {
                    "company": company, "product": product, "link": link, "images": selected_imgs
                }
        else:
            st.info(f"No images found for this type.\n\nTried keys like:\n• Exact: ({_norm(company)}, {_norm(product)}, {_norm(ptype)})\n• Soft / token match variations.\n\nIf using CSV, ensure a row exists with these three columns matching and URLs in ImageURLs.")
else:
    company = st.selectbox("Select Company", sorted(DATA['Company'].dropna().unique()), key="company")
    products = sorted(DATA[DATA['Company'] == company]['Product'].dropna().unique())
    product = st.selectbox("Select Product", products, key="product")
    filtered_rows = DATA[(DATA['Company'] == company) & (DATA['Product'] == product)]

    current_base_key = f"{company}_{product}"
    if st.session_state.last_temp_key and st.session_state.last_temp_key != current_base_key:
        for k, v in st.session_state.temp_selection.items():
            if v['images']:
                st.session_state.ppt_items[f"{v['company']}_{v['product']}_{v['ptype']}"] = v
        st.session_state.temp_selection = {}
    st.session_state.last_temp_key = current_base_key

    for idx, row in filtered_rows.iterrows():
        ptype, link = row['Type'], row.get("Link", "")
        img_paths = [p for p in get_image_list(company, product, ptype) if p and str(p).strip()]

        st.markdown(f"### {ptype}")
        if img_paths:
            img_cols = st.columns(min(4, len(img_paths)))
            sel_key = f"{company}_{product}_{ptype}".replace(" ", "_")
            selected_imgs = st.session_state.temp_selection.get(sel_key, {}).get("images", [])
            for i, path in enumerate(img_paths):
                with img_cols[i % len(img_cols)]:
                    show_image_safe(path)
                    key = f"manual_{company}_{product}_{ptype}_{i}".replace(" ", "_")
                    if st.checkbox("Include", key=key):
                        if path not in selected_imgs:
                            selected_imgs.append(path)
            st.session_state.temp_selection[sel_key] = {
                "company": company, "product": product, "ptype": ptype, "link": link, "images": selected_imgs
            }
        else:
            st.info(
                "No images found for this type.\n\n"
                "Troubleshooting:\n"
                "• If you use a CSV/URL manifest, add a row with columns: Company, Product, Type, ImageURLs\n"
                "• Ensure values match exactly after normalization (case-insensitive, spaces collapsed)\n"
                "• Or place files under images/Company/Product/Type/*.jpg|png"
            )

    for k, v in st.session_state.temp_selection.items():
        if v['images']:
            st.session_state.ppt_items[f"{v['company']}_{v['product']}_{v['ptype']}"] = v

# Sidebar generate PPTs (unchanged)
with st.sidebar:
    st.markdown("## Ready to Download")

    if st.button("📦 Generate Combined PPT"):
        all_items = list(st.session_state.ppt_items.values())
        if all_items:
            prs = create_beautiful_ppt(all_items, include_intro_outro=True)
            ppt_path = "combined_presentation.pptx"
            prs.save(ppt_path)
            st.session_state.ppt_path = ppt_path
            st.session_state.ppt_ready = True
            st.success("PPT generated successfully!")
            st.session_state.ppt_items = {}
            st.session_state.temp_selection = {}
            st.session_state.last_temp_key = None
        else:
            st.warning("No items selected for presentation!")

    if st.session_state.ppt_ready and st.session_state.ppt_path and os.path.exists(st.session_state.ppt_path):
        with open(st.session_state.ppt_path, "rb") as f:
            st.download_button(
                "Download Combined PPT",
                f,
                file_name="combined_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
